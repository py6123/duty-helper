from fastapi import FastAPI, HTTPException, status, UploadFile, File
from fastapi import Form 
import json
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from pymongo import MongoClient
import os
import certifi
from dotenv import load_dotenv
from google import genai

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO
import tempfile
import subprocess
from pathlib import Path

# Load environment variables
load_dotenv()
PROJECT_ID = os.environ.get("GOOGLE_CLOUD_PROJECT")
LOCATION = os.environ.get("GOOGLE_CLOUD_LOCATION", "us-central1")

genai_client = genai.Client(
    vertexai=True,
    project=PROJECT_ID,
    location=LOCATION,
)

app = FastAPI(title="Duty Helper API", description="Direct PyMongo Bridge")

# --- 1. CORS Middleware (Allows your friend's React app to connect) ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"],  
    allow_headers=["*"],  
)

# --- 2. Connect directly to MongoDB ---
# This uses the clean variables from your updated .env file
db = None
collection = None

try:
    #client = MongoClient(os.environ['MONGO_URI'])
    client = MongoClient( 
        os.environ['MONGO_URI'],
        tls=True,
        tlsCAFile=certifi.where()
    )
    db = client[os.environ['MONGO_DB']]
    collection = db[os.environ['MONGO_COLLECTION']]
    # Quick ping to verify connection on startup
    client.admin.command('ping')
    print("Successfully connected to MongoDB!")
except Exception as e:
    print(f"Error connecting to MongoDB: {e}")

# --- 3. Define the Schema ---
class Task(BaseModel):
    task_name: str
    status: str = Field(default="To Do", description="'To Do', 'Doing', or 'Done'")
    eisenhower_category: str
    color_hex: str
    pomodoro_target: int
    pomodoro_completed: int = Field(default=0)
    difficulty: str
    ai_execution_plan: str = Field(default="")
    scheduled_date: str = Field(default="")

class ChatRequest(BaseModel):
    messages: list[dict]
    system: str = ""

# --- 4. The AI Communication Endpoint ---
@app.post("/api/ai/submit-task", status_code=status.HTTP_201_CREATED)
async def ai_submit_task(task: Task):
    if task.pomodoro_target > 10:
         task.pomodoro_target = 10 
         
    task_dict = task.model_dump()
    
    # Insert directly into the collection
    result = collection.insert_one(task_dict)
    
    if result.inserted_id:
        return {"message": "Task successfully saved", "id": str(result.inserted_id)}
    else:
        raise HTTPException(status_code=500, detail="Database insertion failed")

# --- 5. The UI Communication Endpoint ---
@app.get("/api/ui/get-tasks")
async def get_all_tasks():
    tasks = []
    # Fetch all tasks directly from the collection
    for doc in collection.find():
        doc["_id"] = str(doc["_id"]) # Convert MongoDB's special ID to a normal string for React
        tasks.append(doc)
        
    return tasks

@app.post("/api/chat")
async def chat(
    # We use Form() and File() here so React can send text and files in one payload
    message: str = Form(...),
    history: str = Form("[]"), 
    system: str = Form(""),
    files: list[UploadFile] = File(default=[])
):
    # --- 1. The File Bouncer ---
    # Filter out empty file objects just in case React sends an empty array
    valid_files = [f for f in files if f.filename] 
    
    if len(valid_files) > 10:
        raise HTTPException(status_code=400, detail="Maximum 10 files allowed per message.")

    # --- 2. Process Files ONLY when the chat is sent ---
    file_context = ""
    if valid_files:
        for file in valid_files:
            result = await process_single_file(file)
            # Add the extracted text to a hidden context block for Gemini
            file_context += f"\n--- Document: {result['filename']} ---\n{result['text'][:10000]}\n"

    # --- 3. Reconstruct the Conversation ---
    try:
        past_messages = json.loads(history)
    except json.JSONDecodeError:
        past_messages = []

    conversation = "\n".join([
        f"{m.get('role', 'user')}: {m.get('content', '')}"
        for m in past_messages
    ])

    # Assemble the final prompt with the files hidden in the background
    current_input = f"User: {message}\n"
    if file_context:
        current_input += f"\n[Attached Files Content for Context]:\n{file_context}"

    prompt = f"""
System instruction:
{system}

Conversation History:
{conversation}

{current_input}
"""

    # --- 4. Send everything to Vertex AI ---
    response = genai_client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
    )

    return {"text": response.text}

def find_soffice():
    possible_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
    ]

    for path in possible_paths:
        try:
            result = subprocess.run(
                [path, "--version"],
                capture_output=True,
                text=True,
                timeout=10,
                check=False,
            )
            if result.returncode == 0:
                return path
        except Exception:
            pass

    return None


def convert_with_libreoffice(content: bytes, input_ext: str, output_ext: str) -> bytes:
    soffice = find_soffice()
    if not soffice:
        raise HTTPException(
            status_code=400,
            detail=f"{input_ext.upper()} files need LibreOffice installed. Please install LibreOffice or upload {output_ext.upper()} instead.",
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        input_path = tmp_path / f"upload.{input_ext}"
        input_path.write_bytes(content)

        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                output_ext,
                "--outdir",
                str(tmp_path),
                str(input_path),
            ],
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )

        converted_path = tmp_path / f"upload.{output_ext}"

        if result.returncode != 0 or not converted_path.exists():
            raise HTTPException(
                status_code=400,
                detail=f"Could not convert {input_ext.upper()} file. LibreOffice error: {result.stderr or result.stdout}",
            )

        return converted_path.read_bytes()


def extract_pdf_text(content: bytes) -> str:
    pdf = PdfReader(BytesIO(content))
    return "\n".join([
        page.extract_text() or ""
        for page in pdf.pages
    ])


def extract_docx_text(content: bytes) -> str:
    doc = Document(BytesIO(content))
    return "\n".join([
        para.text
        for para in doc.paragraphs
    ])


def extract_pptx_text(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    texts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)

        if slide_text:
            texts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))

    return "\n\n".join(texts)

# ==========================================
# 1. THE HELPER FUNCTION (Shared Logic)
# ==========================================
async def process_single_file(file: UploadFile) -> dict:
    try:
        filename = file.filename.lower()
        content = await file.read()

        if filename.endswith(".pdf"):
            text = extract_pdf_text(content)
        elif filename.endswith(".docx"):
            text = extract_docx_text(content)
        elif filename.endswith(".pptx"):
            text = extract_pptx_text(content)
        elif filename.endswith(".doc"):
            converted = convert_with_libreoffice(content, "doc", "docx")
            text = extract_docx_text(converted)
        elif filename.endswith(".ppt"):
            converted = convert_with_libreoffice(content, "ppt", "pptx")
            text = extract_pptx_text(converted)
        elif filename.endswith((".txt", ".md", ".csv")):
            text = content.decode("utf-8", errors="ignore")
        else:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type. Please upload PDF, DOC, DOCX, PPT, PPTX, TXT, MD, or CSV."
            )

        if not text.strip():
            raise HTTPException(
                status_code=400,
                detail="Could not extract text from this file. It may be scanned image content."
            )

        return {
            "filename": file.filename,
            "text": text
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"File processing error: {str(e)}")

# ==========================================
# 2. BATCH UPLOAD ENDPOINT
# ==========================================
@app.post("/api/upload-notes-batch")
async def upload_notes_batch(files: list[UploadFile] = File(...)):
    
    # --- The Bouncer: Reject anything over 10 files ---
    if len(files) > 10:
        raise HTTPException(
            status_code=400,
            detail=f"Too many files! You uploaded {len(files)}, but the maximum allowed is 10."
        )

    results = []
    
    for file in files:
        try:
            # Calls the Helper Function above
            result = await process_single_file(file) 
            
            # Cap at 10,000 characters per file to save AI tokens
            result["text"] = result["text"][:10000] 
            results.append(result)
            
        except HTTPException as e:
            results.append({"filename": file.filename, "error": e.detail})
        except Exception as e:
            results.append({"filename": file.filename, "error": str(e)})
            
    return {"results": results}

# ==========================================
# 3. SINGLE UPLOAD ENDPOINT
# ==========================================
@app.post("/api/upload-notes")
async def upload_notes(file: UploadFile = File(...)):
    # Calls the exact same Helper Function above
    result = await process_single_file(file)
    
    # Allows a larger character cap for single files
    result["text"] = result["text"][:30000]
    return result

